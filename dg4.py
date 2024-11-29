import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
import io
import shutil
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from dotenv import load_dotenv
import os
from langchain.prompts import PromptTemplate
import re
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor



os.environ['GOOGLE_API_KEY'] = 'AIzaSyAO4CDjz08S2j-79EsPe2uyDtBtNT9TfeY'
load_dotenv()

def get_pdf_text(pdf_docs):
    """Extracts text from uploaded PDF files."""
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text() or ""  # Handle cases where text extraction fails
    return text

def get_text_chunks(text):
    """Splits text into manageable chunks."""
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    """Creates and saves a vector store from text chunks."""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)

    if os.path.exists("faiss_index"):
        shutil.rmtree("faiss_index")
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    """Returns a conversational chain for generating responses."""
    
    # Enhanced Prompt
    prompt_template = """
    You are an intelligent and highly adaptable assistant, capable of understanding and solving a wide variety of tasks. Your primary goal is to assist the user with any type of question or request, including answering complex queries, performing comparisons, distinguishing between concepts, transforming content (e.g., summarization, rephrasing, translation), and providing insightful recommendations.

    When generating content for a PowerPoint presentation:
    - If the user provides headings or subheadings, ensure the content is structured according to these instructions, with clear, relevant bullet points, titles, and summaries. Prioritize organization and clarity, ensuring that the presentation is easy to follow.
    - If no headings or subheadings are provided, analyze the context and content, and create appropriate titles, subtitles, and content sections to ensure the PowerPoint is logical, visually appealing, and contextually relevant.
 

    When answering questions:
    - Provide clear, concise, and correct responses based on the provided PDF or other resources, with relevant context or explanations as needed.
    - When comparing or distinguishing between items, offer a detailed analysis, highlighting similarities, differences, and key factors that the user should consider.

    For content transformation tasks:
    - Modify the content based on the user’s instructions, ensuring it is tailored to their needs. This could include summarization, rephrasing, or reorganization for clarity and impact.
    - If the user requests analysis or logical reasoning, provide a detailed, methodical approach with well-reasoned explanations.

    Flexibility:
    - Be adaptive in your approach, responding to any changes or additional requests efficiently and ensuring that all output is actionable, relevant, and easy to understand, regardless of complexity.

    Unavailable Information:
    - If the necessary information or context is missing or incomplete, clearly state: "The requested information or task cannot be completed based on the provided context."
    
    

    Context:
    {context}

    Question/Task:
    {question}

    Response:
    """
    
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_ppt(response_text):
    """Creates a professional PowerPoint presentation from the generated response."""
    prs = Presentation()

    # Define theme colors
    title_color = RGBColor(0, 0, 0)   # Dark blue for titles
    content_color = RGBColor(0, 0, 0)   # Black for content

    def add_title_slide(prs, title_text, subtitle_text=None):
        """Adds a title slide to the presentation."""
        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = title_color

        # Subtitle
        if subtitle_text:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = content_color

    def add_content_slide(prs, title_text, content_text):
        """Adds a content slide with bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Slide Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = title_color

        # Content Box
        content_box = slide.placeholders[1]
        for line in content_text.strip().split('\n'):
            p = content_box.text_frame.add_paragraph()
            p.text = line.strip()
            p.level = 0  # Top-level bullet
            p.font.size = Pt(18)
            p.font.color.rgb = content_color

    # Parse the response text into structured sections
    lines = response_text.strip().split('\n')
    title_slide_title = lines[0] if len(lines) > 0 else "Generated Report"
    title_slide_subtitle = lines[1] if len(lines) > 1 else None
    add_title_slide(prs, title_slide_title, title_slide_subtitle)

    # Extract slide content
    current_title = None
    current_content = []
    for line in lines[2:]:
        line = line.strip()
        if line.startswith("**Slide"):  # New slide section
            if current_title:
                add_content_slide(prs, current_title, "\n".join(current_content))
            current_title = line.split(':', 1)[-1].strip()  # Get the slide title after ':'
            current_content = []
        else:
            current_content.append(line)

    # Add the last slide
    if current_title:
        add_content_slide(prs, current_title, "\n".join(current_content))

    # Save PowerPoint to a stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    return ppt_stream


def user_input(user_question):
    """Processes user input to generate a response based on the vector store."""
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)

    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question})
    print(response["output_text"])
    return response["output_text"]


def main():
    """Main function to run the Streamlit app."""
    st.set_page_config("Chat PDF")
    st.header("Chat with PDF using Gemini 💁")

    user_question = st.text_input("Ask a Question from the PDF Files")

    if user_question:
        try:
            response_text = user_input(user_question)
            ppt_stream = create_ppt(response_text)
            
            st.success("PowerPoint generated successfully!")
            st.download_button(
                label="Download PowerPoint",
                data=ppt_stream,
                file_name="generated_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"Error: {str(e)}")

    with st.sidebar:
        st.title("Menu:")
        pdf_docs = st.file_uploader("Upload your PDF Files and Click on the Submit & Process Button", accept_multiple_files=True)
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                try:
                    raw_text = get_pdf_text(pdf_docs)
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.success("PDFs processed successfully!")
                except Exception as e:
                    st.error(f"Error: {str(e)}")


        
if __name__ == '__main__':
    main()
